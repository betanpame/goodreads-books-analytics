"""Create a shareable slide deck for Phase 05 · Step 03 results."""
from __future__ import annotations

import argparse
from pathlib import Path
from typing import Sequence

import pandas as pd
from pptx import Presentation
from pptx.util import Inches

DEFAULT_NOTES = Path("docs/phase-05-step-03-task-01-notes.md")
DEFAULT_SUMMARY_CSV = Path("outputs/phase05_step03_task01/comparison_summary.csv")
DEFAULT_CHART = Path("outputs/phase05_step03_task01/comparison_summary.png")
DEFAULT_OUTPUT = Path("docs/phase-05-step-03-task-01-slide.pptx")


def parse_args(argv: Sequence[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Export Phase 05 step 03 summary slides.")
    parser.add_argument("--notes", default=str(DEFAULT_NOTES), help="Path to the task notes markdown file")
    parser.add_argument(
        "--summary-csv",
        default=str(DEFAULT_SUMMARY_CSV),
        help="Path to comparison_summary.csv",
    )
    parser.add_argument(
        "--chart",
        default=str(DEFAULT_CHART),
        help="Path to comparison_summary.png for embedding",
    )
    parser.add_argument(
        "--output",
        default=str(DEFAULT_OUTPUT),
        help="Destination .pptx path",
    )
    return parser.parse_args(argv)


def build_title_slide(prs: Presentation, notes_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Phase 05 · Step 03"
    subtitle = slide.placeholders[1]
    subtitle.text = (
        "SQL vs pandas verification summary\n"
        f"Notes: {notes_path.as_posix()}\n"
        "Automation: scripts/Invoke-Phase05ComparisonRefresh.ps1"
    )


def build_workflow_slide(prs: Presentation) -> None:
    bullets = [
        "Refresh dockerized Python + Postgres stack",
        "Reload curated books_clean table for canonical metrics",
        "Run sql_vs_pandas_compare CLI and capture artifacts",
        "Regenerate comparison_summary.png for visuals",
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Workflow snapshot"
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0


def build_metrics_slide(prs: Presentation, summary_csv: Path) -> None:
    summary = pd.read_csv(summary_csv)
    summary = summary.dropna(subset=["comparison"])
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Metric coverage"
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()
    for idx, row in enumerate(summary.itertuples()):
        bullet = (
            f"{row.comparison}: {int(row.rows_sql)} SQL vs {int(row.rows_pandas)} pandas rows • "
            f"mismatches {int(row.value_mismatches)} ({row.notes})"
        )
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0


def build_chart_slide(prs: Presentation, chart_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = "Visual checkpoint"
    if not chart_path.exists():
        placeholder = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        placeholder.text_frame.text = f"Chart not found at {chart_path}"
        return

    width = Inches(9)
    left = Inches(0.5)
    top = Inches(1.5)
    slide.shapes.add_picture(str(chart_path), left, top, width=width)


def export_slide(notes: Path, summary_csv: Path, chart: Path, output: Path) -> None:
    prs = Presentation()
    build_title_slide(prs, notes)
    build_workflow_slide(prs)
    build_metrics_slide(prs, summary_csv)
    build_chart_slide(prs, chart)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"Saved slide deck to {output}")


def main(argv: Sequence[str] | None = None) -> None:
    args = parse_args(argv)
    export_slide(
        notes=Path(args.notes),
        summary_csv=Path(args.summary_csv),
        chart=Path(args.chart),
        output=Path(args.output),
    )


if __name__ == "__main__":  # pragma: no cover
    main()
